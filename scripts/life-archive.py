#!/usr/bin/env python3
"""
life-archive.py -- raw-first searchable archive for SecondBrain.

This script does not rebuild memory. It creates and maintains a durable local
search index over raw communications and transcripts, and can backfill Gmail
into data/gmail/raw before indexing.
"""

import argparse
import email
import hashlib
import html
import imaplib
import json
import mimetypes
import os
import re
import shutil
import sqlite3
import subprocess
import sys
import tempfile
import time
import zipfile
from concurrent.futures import ProcessPoolExecutor
from datetime import datetime, timezone
from email.header import decode_header
from email.utils import parsedate_to_datetime
from pathlib import Path

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pypdf
except Exception:
    pypdf = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None

try:
    import docx
except Exception:
    docx = None

try:
    import pptx
except Exception:
    pptx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

REPO = Path(__file__).resolve().parents[1]
APPDATA = Path(os.environ.get("APPDATA") or (Path.home() / "AppData" / "Roaming"))
DEFAULT_DB = REPO / "data" / "life-archive" / "life-archive.db"
DEFAULT_GMAIL_RAW = REPO / "data" / "gmail" / "raw"
DEFAULT_GMAIL_BACKFILL_STATE = REPO / "data" / "life-archive" / "gmail-backfill-state.json"
DEFAULT_LOCK = DEFAULT_DB.with_suffix(".lock")
SECRETS = Path.home() / ".secrets"
TESSERACT_CANDIDATES = [
    Path(os.environ.get("TESSERACT_EXE", "")) if os.environ.get("TESSERACT_EXE") else None,
    Path("C:/Program Files/Tesseract-OCR/tesseract.exe"),
]
SMS_ARCHIVE_EXTENSIONS = {".csv", ".txt", ".xlsx", ".pdf", ".vcf", ".url"}
SMS_MEDIA_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".heic", ".mov", ".mp4",
    ".mp3", ".m4a", ".3gp", ".wav", ".amr",
}
MIGRATED_TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".jsonl",
    ".html", ".htm", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".conf", ".log", ".adoc", ".rst", ".sql", ".py", ".js", ".jsx",
    ".ts", ".tsx", ".css", ".scss", ".less", ".vim", ".pl", ".pm",
    ".rb", ".go", ".rs", ".java", ".c", ".cc", ".cpp", ".h", ".hpp",
    ".sh", ".ps1", ".bat", ".cmd", ".properties", ".gradle", ".sln",
    ".csproj", ".vcxproj", ".tex", ".rtf", ".url", ".vcf", ".rels",
    ".svg", ".xsd", ".xslt", ".dtd", ".manifest", ".config",
}
MIGRATED_DOCUMENT_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".xlsm",
}
MIGRATED_IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff",
}
MIGRATED_CONTENT_MAX_BYTES = 16 * 1024 * 1024
MIGRATED_SENSITIVE_PARTS = {
    ".ssh", ".gnupg", "session storage", "sessions", "cookies", "local storage",
    "indexeddb", "password stores", "password-store", "credential stores",
}
MIGRATED_SENSITIVE_NAMES = {
    ".env", "login data", "web data", "cookies", "credentials", "secrets.json",
    "key4.db", "logins.json", "signons.sqlite", "wallet.dat", "id_rsa", "id_dsa",
    "id_ecdsa", "id_ed25519", ".npmrc", ".pypirc", ".netrc", "auth.json",
    "service-account.json",
}
MIGRATED_PRIVATE_KEY_EXTENSIONS = {".pem", ".p12", ".pfx", ".key", ".ppk", ".kdbx"}
MIGRATED_NOISY_PARTS = {
    "node_modules", ".git", "vendor", "dist", "build", "target", "coverage",
    "__pycache__", ".cache", "cache", "temp", "tmp",
}
MIGRATED_CODE_EXTENSIONS = {
    ".py", ".js", ".jsx", ".ts", ".tsx", ".css", ".scss", ".less", ".vim",
    ".pl", ".pm", ".rb", ".go", ".rs", ".java", ".c", ".cc", ".cpp",
    ".h", ".hpp", ".sh", ".ps1", ".bat", ".cmd", ".dll", ".exe",
}


def now_iso():
    return datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")


def sha256_bytes(data):
    return hashlib.sha256(data).hexdigest()


def sha256_file(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def safe_segment(value, fallback="item", max_len=120):
    value = str(value or fallback).strip()
    value = re.sub(r"[^A-Za-z0-9._@+-]+", "-", value)
    value = value.strip(".-") or fallback
    return value[:max_len]


def decode_hdr(value):
    if not value:
        return ""
    out = []
    for part, enc in decode_header(value):
        if isinstance(part, bytes):
            try:
                out.append(part.decode(enc or "utf-8", errors="replace"))
            except LookupError:
                out.append(part.decode("utf-8", errors="replace"))
        else:
            out.append(str(part))
    return "".join(out)


def parse_email_date(value):
    try:
        dt = parsedate_to_datetime(value or "")
        if dt.tzinfo is None:
            dt = dt.replace(tzinfo=timezone.utc)
        return dt.astimezone(timezone.utc)
    except Exception:
        return None


def received_date_header(msg):
    for header in reversed(msg.get_all("Received", []) or []):
        if ";" not in header:
            continue
        candidate = header.rsplit(";", 1)[-1].strip()
        if parse_email_date(candidate) is not None:
            return candidate
    return ""


def date_parts(value):
    dt = None
    if isinstance(value, datetime):
        dt = value
    elif value:
        try:
            dt = datetime.fromisoformat(str(value).replace("Z", "+00:00"))
        except Exception:
            try:
                dt = parsedate_to_datetime(str(value))
            except Exception:
                dt = None
    if dt is None:
        dt = datetime.now(timezone.utc)
    return f"{dt.year:04d}", f"{dt.month:02d}", f"{dt.day:02d}"


def utf8_safe(value):
    return str(value or "").encode("utf-8", errors="replace").decode("utf-8", errors="replace")


def safe_json_dumps(value, **kwargs):
    kwargs.setdefault("ensure_ascii", False)
    return json.dumps(value, **kwargs).encode("utf-8", errors="replace").decode("utf-8", errors="replace")


def clean_text(value):
    text = utf8_safe(value)
    text = re.sub(r"<(script|style)\b[^>]*>[\s\S]*?</\1>", " ", text, flags=re.I)
    text = re.sub(r"<[^>]+>", " ", html.unescape(text))
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def ensure_parent(path):
    Path(path).parent.mkdir(parents=True, exist_ok=True)


def connect_db(db_path=DEFAULT_DB):
    db_path = Path(db_path)
    db_path.parent.mkdir(parents=True, exist_ok=True)
    con = sqlite3.connect(str(db_path))
    con.row_factory = sqlite3.Row
    con.execute("PRAGMA busy_timeout=30000")
    con.execute("PRAGMA journal_mode=WAL")
    con.execute("PRAGMA synchronous=NORMAL")
    con.execute(
        """
        CREATE TABLE IF NOT EXISTS items (
          id TEXT PRIMARY KEY,
          source TEXT NOT NULL,
          source_id TEXT,
          kind TEXT,
          title TEXT,
          author TEXT,
          recipients TEXT,
          created_at TEXT,
          raw_path TEXT,
          url TEXT,
          checksum TEXT,
          metadata_json TEXT,
          indexed_at TEXT NOT NULL
        )
        """
    )
    con.execute(
        """
        CREATE TABLE IF NOT EXISTS attachments (
          id TEXT PRIMARY KEY,
          item_id TEXT NOT NULL,
          filename TEXT,
          content_type TEXT,
          path TEXT,
          size INTEGER,
          checksum TEXT,
          extracted_text TEXT,
          metadata_json TEXT,
          FOREIGN KEY(item_id) REFERENCES items(id)
        )
        """
    )
    con.execute(
        """
        CREATE VIRTUAL TABLE IF NOT EXISTS item_fts USING fts5(
          item_id UNINDEXED,
          source,
          title,
          body,
          attachment_text,
          participants,
          raw_path UNINDEXED,
          tokenize='porter unicode61'
        )
        """
    )
    con.execute("CREATE INDEX IF NOT EXISTS idx_items_source ON items(source, created_at)")
    con.execute("CREATE INDEX IF NOT EXISTS idx_items_raw_path ON items(raw_path)")
    return con


def process_is_running(pid):
    try:
        pid = int(pid)
    except Exception:
        return False
    if pid <= 0:
        return False
    if os.name == "nt":
        try:
            result = subprocess.run(
                ["tasklist", "/FI", f"PID eq {pid}", "/FO", "CSV", "/NH"],
                capture_output=True,
                text=True,
                timeout=5,
            )
            return f'"{pid}"' in (result.stdout or "")
        except Exception:
            return True
    try:
        os.kill(pid, 0)
        return True
    except OSError:
        return False


def acquire_lock(path=DEFAULT_LOCK, stale_seconds=1800):
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    now = time.time()
    if path.exists():
        lock_pid = None
        try:
            lock_pid = (json.loads(path.read_text(encoding="utf-8", errors="replace")) or {}).get("pid")
        except Exception:
            lock_pid = None
        if lock_pid and not process_is_running(lock_pid):
            try:
                path.unlink()
            except OSError:
                return None
        if not path.exists():
            pass
        else:
            try:
                age = now - path.stat().st_mtime
            except OSError:
                age = 0
            if age < stale_seconds:
                return None
            try:
                path.unlink()
            except OSError:
                return None
    if path.exists():
        try:
            path.unlink()
        except OSError:
            return None
    try:
        fd = os.open(str(path), os.O_CREAT | os.O_EXCL | os.O_WRONLY)
        os.write(fd, safe_json_dumps({"pid": os.getpid(), "createdAt": now_iso()}).encode("utf-8"))
        return fd
    except FileExistsError:
        return None


def release_lock(fd, path=DEFAULT_LOCK):
    if fd is None:
        return
    try:
        os.close(fd)
    except OSError:
        pass
    try:
        Path(path).unlink()
    except OSError:
        pass


def lock_path_for_db(db_path):
    try:
        return Path(db_path).with_suffix(".lock")
    except Exception:
        return DEFAULT_LOCK


def item_id_for(source, source_id, raw_path):
    basis = f"{source}|{source_id or ''}|{Path(raw_path).as_posix() if raw_path else ''}"
    return hashlib.sha1(basis.encode("utf-8", errors="replace")).hexdigest()[:32]


def item_checksum_matches(con, item_id, checksum):
    if not checksum:
        return False
    row = con.execute("SELECT checksum FROM items WHERE id = ?", (item_id,)).fetchone()
    return bool(row and row["checksum"] == checksum)


def upsert_item(con, item, body="", attachments=None):
    attachments = attachments or []
    item = dict(item)
    item.setdefault("id", item_id_for(item.get("source"), item.get("source_id"), item.get("raw_path")))
    item.setdefault("indexed_at", now_iso())
    item.setdefault("metadata_json", "{}")
    fields = [
        "id", "source", "source_id", "kind", "title", "author", "recipients",
        "created_at", "raw_path", "url", "checksum", "metadata_json", "indexed_at",
    ]
    for field in fields:
        if isinstance(item.get(field), str):
            item[field] = utf8_safe(item[field])
    con.execute(
        """
        INSERT INTO items ({fields}) VALUES ({marks})
        ON CONFLICT(id) DO UPDATE SET
          source=excluded.source,
          source_id=excluded.source_id,
          kind=excluded.kind,
          title=excluded.title,
          author=excluded.author,
          recipients=excluded.recipients,
          created_at=excluded.created_at,
          raw_path=excluded.raw_path,
          url=excluded.url,
          checksum=excluded.checksum,
          metadata_json=excluded.metadata_json,
          indexed_at=excluded.indexed_at
        """.format(fields=", ".join(fields), marks=", ".join(["?"] * len(fields))),
        [item.get(f) for f in fields],
    )
    con.execute("DELETE FROM item_fts WHERE item_id = ?", (item["id"],))
    body = utf8_safe(body or "")
    attachment_text = "\n".join(utf8_safe(a.get("extracted_text") or "") for a in attachments)
    participants = utf8_safe(" ".join([item.get("author") or "", item.get("recipients") or ""]))
    con.execute(
        """
        INSERT INTO item_fts (item_id, source, title, body, attachment_text, participants, raw_path)
        VALUES (?, ?, ?, ?, ?, ?, ?)
        """,
        (
            item["id"],
            item.get("source") or "",
            item.get("title") or "",
            body or "",
            attachment_text,
            participants,
            item.get("raw_path") or "",
        ),
    )
    con.execute("DELETE FROM attachments WHERE item_id = ?", (item["id"],))
    for att_idx, att in enumerate(attachments):
        att_id = hashlib.sha1(
            f"{item['id']}|{att_idx}|{att.get('path') or att.get('filename')}".encode("utf-8", errors="replace")
        ).hexdigest()[:32]
        con.execute(
            """
            INSERT INTO attachments
            (id, item_id, filename, content_type, path, size, checksum, extracted_text, metadata_json)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                att_id,
                item["id"],
                utf8_safe(att.get("filename")),
                utf8_safe(att.get("content_type")),
                utf8_safe(att.get("path")),
                att.get("size"),
                utf8_safe(att.get("checksum")),
                utf8_safe(att.get("extracted_text")),
                safe_json_dumps(att.get("metadata") or {}),
            ),
        )
    return item["id"]


def tesseract_path():
    for p in TESSERACT_CANDIDATES:
        if p and p.exists():
            return str(p)
    found = shutil.which("tesseract")
    return found


def extract_text_from_image(path):
    exe = tesseract_path()
    if not exe:
        return ""
    try:
        with tempfile.TemporaryDirectory() as td:
            out_base = Path(td) / "ocr"
            subprocess.run(
                [exe, str(path), str(out_base), "--psm", "6"],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                timeout=60,
                check=False,
            )
            out = out_base.with_suffix(".txt")
            if out.exists():
                return clean_text(out.read_text(encoding="utf-8", errors="replace"))
    except Exception:
        return ""
    return ""


def extract_text_from_pdf(path):
    texts = []
    if pdfplumber:
        try:
            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages[:50]:
                    texts.append(page.extract_text() or "")
        except Exception:
            pass
    if not any(t.strip() for t in texts) and pypdf:
        try:
            reader = pypdf.PdfReader(str(path))
            for page in reader.pages[:50]:
                texts.append(page.extract_text() or "")
        except Exception:
            pass
    return clean_text("\n".join(texts))


def extract_text_from_docx(path):
    if not docx:
        return ""
    try:
        d = docx.Document(str(path))
        return clean_text("\n".join(p.text for p in d.ExampleCoraphs))
    except Exception:
        return ""


def extract_text_from_pptx(path):
    if not pptx:
        return ""
    try:
        deck = pptx.Presentation(str(path))
        texts = []
        for slide in deck.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return clean_text("\n".join(texts))
    except Exception:
        return ""


def extract_text_from_xlsx(path):
    if not openpyxl:
        return ""
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        values = []
        for ws in wb.worksheets[:10]:
            values.append(ws.title)
            for row in ws.iter_rows(max_row=300, values_only=True):
                vals = [str(v) for v in row if v is not None]
                if vals:
                    values.append(" | ".join(vals))
        return clean_text("\n".join(values))
    except Exception:
        return ""


def extract_text_from_zip_xml(path):
    try:
        with zipfile.ZipFile(path) as z:
            texts = []
            for name in z.namelist():
                if not name.endswith(".xml"):
                    continue
                if not any(part in name for part in ("word/", "ppt/", "xl/")):
                    continue
                raw = z.read(name).decode("utf-8", errors="ignore")
                texts.append(clean_text(raw))
            return clean_text("\n".join(texts))
    except Exception:
        return ""


def extract_attachment_text(path, content_type="", write_sidecar=True):
    path = Path(path)
    ext = path.suffix.lower()
    text = ""
    if ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff") or content_type.startswith("image/"):
        text = extract_text_from_image(path)
    elif ext == ".pdf" or content_type == "application/pdf":
        text = extract_text_from_pdf(path)
    elif ext == ".docx":
        text = extract_text_from_docx(path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(path)
    elif ext in (".xlsx", ".xlsm"):
        text = extract_text_from_xlsx(path)
    elif ext in (".txt", ".md", ".csv", ".json", ".html", ".htm", ".log"):
        try:
            text = path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            text = ""
    if not text and ext in (".docx", ".pptx", ".xlsx", ".xlsm"):
        text = extract_text_from_zip_xml(path)
    text = clean_text(text)
    if text and write_sidecar:
        sidecar = path.with_suffix(path.suffix + ".text.txt")
        try:
            sidecar.write_text(text, encoding="utf-8")
        except Exception:
            pass
    return text


def iso_utc(value):
    if not value:
        return None
    if isinstance(value, str):
        raw = value.strip()
        if raw.startswith("D:") and len(raw) >= 16:
            raw = raw[2:16]
            try:
                return datetime.strptime(raw, "%Y%m%d%H%M%S").replace(tzinfo=timezone.utc).isoformat()
            except Exception:
                return None
        try:
            value = datetime.fromisoformat(raw.replace("Z", "+00:00"))
        except Exception:
            return None
    if not isinstance(value, datetime):
        return None
    if value.tzinfo is None:
        value = value.replace(tzinfo=timezone.utc)
    return value.astimezone(timezone.utc).isoformat()


def filename_date(path):
    stem = Path(path).stem
    patterns = [
        (r"(?<!\d)((?:19|20)\d{2})[-_. ]([01]\d)[-_. ]([0-3]\d)(?!\d)", "%Y%m%d"),
        (r"(?<!\d)((?:19|20)\d{2})([01]\d)([0-3]\d)(?!\d)", "%Y%m%d"),
        (r"(?<!\d)([01]\d)[-_. ]([0-3]\d)[-_. ]((?:19|20)\d{2})(?!\d)", "%m%d%Y"),
    ]
    for pattern, fmt in patterns:
        match = re.search(pattern, stem)
        if not match:
            continue
        try:
            dt = datetime.strptime("".join(match.groups()), fmt).replace(hour=12, tzinfo=timezone.utc)
            return dt.isoformat()
        except ValueError:
            continue
    return None


def embedded_file_times(path):
    path = Path(path)
    ext = path.suffix.lower()
    created = None
    modified = None
    try:
        if ext == ".docx" and docx:
            props = docx.Document(str(path)).core_properties
            created, modified = props.created, props.modified
        elif ext == ".pptx" and pptx:
            props = pptx.Presentation(str(path)).core_properties
            created, modified = props.created, props.modified
        elif ext in (".xlsx", ".xlsm") and openpyxl:
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            created, modified = wb.properties.created, wb.properties.modified
            wb.close()
        elif ext == ".pdf" and pypdf:
            meta = pypdf.PdfReader(str(path)).metadata or {}
            created = getattr(meta, "creation_date", None) or meta.get("/CreationDate")
            modified = getattr(meta, "modification_date", None) or meta.get("/ModDate")
        elif ext in MIGRATED_IMAGE_EXTENSIONS and Image:
            with Image.open(str(path)) as image:
                exif = image.getexif()
                raw_created = exif.get(36867) or exif.get(306)
                if raw_created:
                    created = datetime.strptime(str(raw_created), "%Y:%m:%d %H:%M:%S")
    except Exception:
        pass
    return iso_utc(created), iso_utc(modified)


def migrated_file_sensitivity(relative_path):
    rel = Path(relative_path)
    parts = [part.lower() for part in rel.parts]
    name = rel.name.lower()
    ext = rel.suffix.lower()
    if any(part in MIGRATED_SENSITIVE_PARTS for part in parts):
        return "secret_or_session_store"
    if name in MIGRATED_SENSITIVE_NAMES or name.startswith(".env."):
        return "credential_or_secret_file"
    if re.search(r"(?:credential|secret|password|access[-_]?token|refresh[-_]?token)", name):
        return "credential_or_secret_file"
    if ext in MIGRATED_PRIVATE_KEY_EXTENSIONS:
        return "private_key_or_encrypted_vault"
    if re.fullmatch(r"id_(rsa|dsa|ecdsa|ed25519)(\.pub)?", name) and not name.endswith(".pub"):
        return "private_key"
    return None


def migrated_file_is_noisy(relative_path):
    return any(part.lower() in MIGRATED_NOISY_PARTS for part in Path(relative_path).parts)


def migrated_file_looks_textual(path):
    try:
        sample = Path(path).read_bytes()[:8192]
    except Exception:
        return False
    if not sample or b"\x00" in sample:
        return False
    printable = sum(1 for byte in sample if byte in b"\t\n\r" or 32 <= byte <= 126 or byte >= 128)
    return printable / len(sample) >= 0.90


def migrated_graphiti_eligible(relative_path, extension, extraction_status, sensitivity):
    rel = Path(relative_path)
    parts = [part.lower() for part in rel.parts]
    if sensitivity or extraction_status != "text_extracted":
        return False
    if extension in MIGRATED_CODE_EXTENSIONS or migrated_file_is_noisy(rel):
        return False
    if any(part in {"taxes", "recruiting", "resumes", "personal", "_job applications"} for part in parts):
        return False
    return extension in (MIGRATED_DOCUMENT_EXTENSIONS | MIGRATED_TEXT_EXTENSIONS)


def migrated_timestamp_metadata(path):
    stat = Path(path).stat()
    fs_created = datetime.fromtimestamp(stat.st_ctime, timezone.utc).isoformat()
    fs_modified = datetime.fromtimestamp(stat.st_mtime, timezone.utc).isoformat()
    embedded_created, embedded_modified = embedded_file_times(path)
    dated_name = filename_date(path)
    if embedded_modified:
        best, provenance = embedded_modified, "embedded_modified"
    elif embedded_created:
        best, provenance = embedded_created, "embedded_created"
    elif dated_name:
        best, provenance = dated_name, "filename_date"
    else:
        best, provenance = fs_modified, "filesystem_modified_after_aws_transfer"
    return {
        "best_timestamp": best,
        "timestamp_provenance": provenance,
        "filesystem_created": fs_created,
        "filesystem_modified": fs_modified,
        "filesystem_mtime_ns": stat.st_mtime_ns,
        "embedded_created": embedded_created,
        "embedded_modified": embedded_modified,
        "filename_date": dated_name,
    }


def prepare_migrated_file(path, source_root):
    path = Path(path).resolve()
    source_root = Path(source_root).resolve()
    relative_path = path.relative_to(source_root).as_posix()
    stat = path.stat()
    extension = path.suffix.lower()
    content_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    sensitivity = migrated_file_sensitivity(relative_path)
    noisy = migrated_file_is_noisy(relative_path)
    checksum = sha256_file(path)
    source_id = relative_path
    item_id = item_id_for("migrated-file", source_id, str(path))

    extraction_status = "metadata_only_unsupported"
    extracted = ""
    if sensitivity:
        extraction_status = "metadata_only_sensitive"
    elif noisy:
        extraction_status = "metadata_only_vendor_or_generated"
    elif stat.st_size > MIGRATED_CONTENT_MAX_BYTES:
        extraction_status = "metadata_only_size_limit"
    elif extension in MIGRATED_TEXT_EXTENSIONS or (not extension and migrated_file_looks_textual(path)):
        try:
            extracted = path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            extracted = ""
        extraction_status = "text_extracted" if extracted.strip() else "text_empty_or_unreadable"
    elif extension in (MIGRATED_DOCUMENT_EXTENSIONS | MIGRATED_IMAGE_EXTENSIONS):
        extracted = extract_attachment_text(path, content_type, write_sidecar=False)
        extraction_status = "text_extracted" if extracted else (
            "ocr_unavailable" if extension in MIGRATED_IMAGE_EXTENSIONS and not tesseract_path()
            else "text_empty_or_unreadable"
        )

    times = migrated_timestamp_metadata(path)
    graphiti_eligible = migrated_graphiti_eligible(
        relative_path, extension, extraction_status, sensitivity
    )
    catalog = (
        f"Migrated file catalog entry. Relative path: {relative_path}. "
        f"Filename: {path.name}. Extension: {extension or '[none]'}. MIME: {content_type}."
    )
    body = catalog + (("\n\n" + extracted) if extracted else "")
    metadata = {
        "absolute_path": str(path),
        "source_relative_path": relative_path,
        "migration_root": str(source_root),
        "filename": path.name,
        "extension": extension,
        "mime_type": content_type,
        "size": stat.st_size,
        "sha256": checksum,
        **times,
        "extraction_status": extraction_status,
        "extracted_text_chars": len(clean_text(extracted)) if extracted else 0,
        "sensitivity": sensitivity,
        "vendor_or_generated": noisy,
        "graphiti_eligible": graphiti_eligible,
        "original_timestamp_note": (
            "AWS S3 sync did not retain source filesystem timestamps; embedded document/image "
            "metadata or a filename date is preferred when present."
        ),
    }
    item = {
        "id": item_id,
        "source": "migrated-file",
        "source_id": source_id,
        "kind": "migrated-file",
        "title": path.name,
        "author": "",
        "recipients": "ExampleCo",
        "created_at": times["best_timestamp"],
        "raw_path": str(path),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(metadata),
    }
    return {"item": item, "body": body, "metadata": metadata}


def prepare_migrated_file_safe(job):
    path, source_root = job
    try:
        return {"ok": True, "prepared": prepare_migrated_file(path, source_root)}
    except Exception as exc:
        return {"ok": False, "path": str(path), "error": str(exc)[:500]}


def store_prepared_migrated_file(con, prepared):
    item = prepared["item"]
    if item_checksum_matches(con, item["id"], item.get("checksum")):
        return 0
    upsert_item(con, item, prepared.get("body") or "", [])
    return 1


def index_migrated_file(con, path, source_root):
    return store_prepared_migrated_file(con, prepare_migrated_file(path, source_root))


def html_to_text(value):
    return clean_text(value)


def body_from_email_message(msg):
    plain = []
    html_parts = []
    for part in msg.walk():
        ctype = part.get_content_type()
        disp = str(part.get("Content-Disposition") or "").lower()
        if "attachment" in disp:
            continue
        if ctype not in ("text/plain", "text/html"):
            continue
        payload = part.get_payload(decode=True) or b""
        try:
            text = payload.decode(part.get_content_charset() or "utf-8", errors="replace")
        except LookupError:
            text = payload.decode("utf-8", errors="replace")
        if ctype == "text/html":
            html_parts.append(html_to_text(text))
        else:
            plain.append(clean_text(text))
    return "\n".join([p for p in plain if p] or [p for p in html_parts if p])


def index_gmail_message_json(con, path):
    data = json.loads(Path(path).read_text(encoding="utf-8", errors="replace"))
    raw_path = Path(data.get("raw_saved_to") or Path(path).with_name("raw.eml"))
    checksum = sha256_file(raw_path) if raw_path.exists() else sha256_file(path)
    source_id = data.get("message_id") or data.get("id") or str(path)
    raw_for_id = str(raw_path if raw_path.exists() else path)
    item_id = item_id_for("gmail", source_id, raw_for_id)
    if item_checksum_matches(con, item_id, checksum):
        return 0
    body = clean_text(data.get("body") or data.get("text") or data.get("snippet") or "")
    if not body and raw_path.exists():
        try:
            body = body_from_email_message(email.message_from_bytes(raw_path.read_bytes()))
        except Exception:
            body = ""
    attachments = []
    for att in data.get("attachments") or []:
        att_path = att.get("archive_saved_to") or att.get("saved_to")
        extracted = ""
        checksum_att = None
        size = att.get("size")
        if att_path and Path(att_path).exists():
            checksum_att = sha256_file(att_path)
            extracted = extract_attachment_text(att_path, att.get("content_type") or "")
            size = Path(att_path).stat().st_size
        attachments.append({
            "filename": att.get("filename"),
            "content_type": att.get("content_type"),
            "path": att_path,
            "size": size,
            "checksum": checksum_att,
            "extracted_text": extracted,
            "metadata": att,
        })
    item = {
        "source": "gmail",
        "source_id": source_id,
        "kind": "email",
        "title": data.get("subject") or "(no subject)",
        "author": data.get("from"),
        "recipients": data.get("to"),
        "created_at": data.get("date"),
        "raw_path": raw_for_id,
        "url": data.get("gmail_url"),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(data),
    }
    return upsert_item(con, item, body, attachments)


def index_eml(con, path, source="gmail"):
    raw = Path(path).read_bytes()
    msg = email.message_from_bytes(raw)
    source_id = (msg.get("Message-ID") or "").strip("<>") or sha256_bytes(raw)[:24]
    checksum = sha256_bytes(raw)
    if item_checksum_matches(con, item_id_for(source, source_id, str(path)), checksum):
        return 0
    body = body_from_email_message(msg)
    item = {
        "source": source,
        "source_id": source_id,
        "kind": "email",
        "title": decode_hdr(msg.get("Subject")) or "(no subject)",
        "author": decode_hdr(msg.get("From")),
        "recipients": decode_hdr(msg.get("To")),
        "created_at": decode_hdr(msg.get("Date")),
        "raw_path": str(path),
        "url": None,
        "checksum": checksum,
        "metadata_json": safe_json_dumps({"message_id": source_id}),
    }
    return upsert_item(con, item, body, [])


def index_otter_json(con, path):
    data = json.loads(Path(path).read_text(encoding="utf-8", errors="replace"))
    checksum = sha256_file(path)
    source_id = data.get("id") or data.get("otid") or data.get("speech_id") or str(path)
    if item_checksum_matches(con, item_id_for("otter", source_id, str(path)), checksum):
        return 0
    transcript = data.get("transcript") or data.get("body") or data.get("summary") or ""
    # Otter exports often bury utterance text in nested segment fields.
    nested_text = clean_text(" ".join(flatten_json_text(data)))
    if nested_text and nested_text not in transcript:
        transcript = "\n".join([transcript, nested_text]).strip()
    item = {
        "source": "otter",
        "source_id": source_id,
        "kind": "transcript",
        "title": data.get("title") or "Otter transcript",
        "author": data.get("speaker") or "",
        "recipients": "",
        "created_at": data.get("created_at") or data.get("start_time") or data.get("ingested_at"),
        "raw_path": str(path),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(data),
    }
    return upsert_item(con, item, transcript, [])


def index_vapi_json(con, path):
    data = json.loads(Path(path).read_text(encoding="utf-8", errors="replace"))
    checksum = sha256_file(path)
    msg = data.get("msg") or data
    call_obj = data.get("callObj") or msg.get("call") or {}
    call_id = (call_obj.get("id") if isinstance(call_obj, dict) else None) or msg.get("callId") or Path(path).stem
    if item_checksum_matches(con, item_id_for("vapi", call_id, str(path)), checksum):
        return 0
    transcript = msg.get("transcript") or data.get("transcript") or ""
    summary = (msg.get("analysis") or {}).get("summary") if isinstance(msg.get("analysis"), dict) else ""
    body = "\n\n".join(x for x in [summary, transcript] if x)
    item = {
        "source": "vapi",
        "source_id": call_id,
        "kind": "amy-call",
        "title": f"Amy call: {call_id}",
        "author": str((call_obj.get("customer") or {}).get("number") if isinstance(call_obj, dict) else ""),
        "recipients": "Amy",
        "created_at": data.get("savedAt") or data.get("created_at") or msg.get("startedAt"),
        "raw_path": str(path),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(data),
    }
    return upsert_item(con, item, body, [])


def index_linkedin_json(con, path):
    data = json.loads(Path(path).read_text(encoding="utf-8", errors="replace"))
    checksum = sha256_file(path)
    if isinstance(data.get("threads"), list):
        count = 0
        for i, thread in enumerate(data["threads"]):
            source_id = f"{path}:{i}:{thread.get('name') or thread.get('contact') or ''}"
            if item_checksum_matches(con, item_id_for("linkedin-dm", source_id, str(path)), checksum):
                continue
            item = {
                "source": "linkedin-dm",
                "source_id": source_id,
                "kind": "linkedin-message",
                "title": thread.get("name") or "LinkedIn DM",
                "author": thread.get("name") or "",
                "recipients": "ExampleCo",
                "created_at": data.get("lastCrawl") or thread.get("time"),
                "raw_path": str(path),
                "checksum": checksum,
                "metadata_json": safe_json_dumps(thread),
            }
            upsert_item(con, item, thread.get("snippet") or thread.get("body") or thread.get("text") or "", [])
            count += 1
        return count
    contact = data.get("contact") or {}
    body = data.get("body") or "\n".join(
        p.get("text", "") if isinstance(p, dict) else str(p)
        for p in (data.get("posts") or data.get("recentPosts") or [])
    )
    source_id = data.get("id") or data.get("url") or str(path)
    if item_checksum_matches(con, item_id_for("linkedin", source_id, str(path)), checksum):
        return 0
    item = {
        "source": "linkedin",
        "source_id": source_id,
        "kind": "linkedin-posts",
        "title": f"LinkedIn: {contact.get('name') or data.get('name') or Path(path).stem}",
        "author": contact.get("name") or data.get("name"),
        "recipients": "ExampleCo network",
        "created_at": data.get("scanned_at") or data.get("lastCrawlAt"),
        "raw_path": str(path),
        "url": data.get("url") or contact.get("url"),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(data),
    }
    return upsert_item(con, item, body, [])


def index_jsonl(con, path, source):
    count = 0
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        for n, line in enumerate(f, 1):
            line = line.strip()
            if not line:
                continue
            try:
                data = json.loads(line)
            except Exception:
                continue
            body = data.get("comment") or data.get("text") or data.get("prompt") or data.get("action_summary") or safe_json_dumps(data)
            item = {
                "source": source,
                "source_id": data.get("id") or data.get("gmail_message_id") or data.get("thread_id") or f"{path}:{n}",
                "kind": data.get("source") or data.get("kind") or "jsonl",
                "title": data.get("title") or data.get("subject") or f"{source} row",
                "author": data.get("from") or data.get("origin") or data.get("source"),
                "recipients": data.get("to") or "",
                "created_at": data.get("ts") or data.get("createdAt") or data.get("dispatched_at"),
                "raw_path": str(path),
                "checksum": sha256_file(path),
                "metadata_json": safe_json_dumps(data),
            }
            upsert_item(con, item, body, [])
            count += 1
    return count


def extract_session_text_from_jsonl(path):
    turns = []
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                obj = json.loads(line)
            except Exception:
                continue
            role = obj.get("role") or obj.get("type") or obj.get("speaker") or ""
            text = ""
            msg = obj.get("message")
            if isinstance(msg, dict):
                content = msg.get("content")
                if isinstance(content, str):
                    text = content
                elif isinstance(content, list):
                    bits = []
                    for c in content:
                        if isinstance(c, dict):
                            bits.append(str(c.get("text") or c.get("content") or ""))
                        else:
                            bits.append(str(c))
                    text = "\n".join(bits)
            if not text:
                content = obj.get("content")
                if isinstance(content, str):
                    text = content
                elif isinstance(content, list):
                    text = "\n".join(str(x.get("text") if isinstance(x, dict) else x) for x in content)
            if text:
                turns.append(f"{role}: {clean_text(text)}")
    return "\n".join(turns)


def archive_external_file(path, source, archive_root=None):
    path = Path(path)
    archive_root = Path(archive_root or (REPO / "data" / "sessions" / "raw"))
    try:
        created = datetime.fromtimestamp(path.stat().st_mtime, timezone.utc)
    except Exception:
        created = datetime.now(timezone.utc)
    y, m, d = date_parts(created)
    checksum = sha256_file(path)
    dest_dir = archive_root / safe_segment(source) / y / m / d
    dest_dir.mkdir(parents=True, exist_ok=True)
    dest = dest_dir / f"{checksum[:16]}-{safe_segment(path.name)}"
    if not dest.exists():
        shutil.copy2(path, dest)
    return dest


def index_session_file(con, path, source="session", copy_external=False):
    indexed_path = Path(path)
    if copy_external and not str(indexed_path.resolve()).lower().startswith(str(REPO.resolve()).lower()):
        indexed_path = archive_external_file(indexed_path, source)
    checksum = sha256_file(indexed_path)
    item_id = item_id_for(source, str(indexed_path), str(indexed_path))
    if item_checksum_matches(con, item_id, checksum):
        return None
    suffix = indexed_path.suffix.lower()
    if suffix == ".jsonl":
        text = extract_session_text_from_jsonl(indexed_path)
    elif suffix in (".json", ".txt", ".md", ".log"):
        text = indexed_path.read_text(encoding="utf-8", errors="replace")
    else:
        text = extract_attachment_text(indexed_path)
    item = {
        "source": source,
        "id": item_id,
        "source_id": str(indexed_path),
        "kind": "prompt-session",
        "title": f"{source}: {indexed_path.name}",
        "author": "ExampleCo / assistant",
        "recipients": "",
        "created_at": datetime.fromtimestamp(indexed_path.stat().st_mtime, timezone.utc).isoformat(),
        "raw_path": str(indexed_path),
        "checksum": checksum,
        "metadata_json": safe_json_dumps({"original_path": str(path)}),
    }
    return upsert_item(con, item, text, [])


def flatten_json_text(value, out=None):
    out = out if out is not None else []
    if isinstance(value, dict):
        for k, v in value.items():
            out.append(str(k))
            flatten_json_text(v, out)
    elif isinstance(value, list):
        for v in value:
            flatten_json_text(v, out)
    elif value is not None:
        out.append(str(value))
    return out


def index_structured_raw_file(con, path, source, kind):
    path = Path(path)
    checksum = sha256_file(path)
    item_id = item_id_for(source, str(path), str(path))
    if item_checksum_matches(con, item_id, checksum):
        return None
    text = ""
    metadata = {"original_path": str(path)}
    if path.suffix.lower() == ".json":
        try:
            data = json.loads(path.read_text(encoding="utf-8", errors="replace"))
            metadata["json_keys"] = list(data.keys())[:50] if isinstance(data, dict) else []
            text = clean_text(" ".join(flatten_json_text(data)))
        except Exception:
            text = path.read_text(encoding="utf-8", errors="replace")
    else:
        text = path.read_text(encoding="utf-8", errors="replace")
    item = {
        "source": source,
        "id": item_id,
        "source_id": str(path),
        "kind": kind,
        "title": f"{source}: {path.name}",
        "author": "",
        "recipients": "ExampleCo",
        "created_at": datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
        "raw_path": str(path),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(metadata),
    }
    return upsert_item(con, item, text, [])


def index_whatsapp_conversation(con, path):
    path = Path(path)
    checksum = sha256_file(path)
    if item_checksum_matches(con, item_id_for("whatsapp", str(path), str(path)), checksum):
        return 0
    transcript = path.read_text(encoding="utf-8", errors="replace")
    meta_path = path.with_name("meta.json")
    meta = {}
    if meta_path.exists():
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8", errors="replace"))
        except Exception:
            meta = {}
    first_line = transcript.splitlines()[0] if transcript.splitlines() else "WhatsApp conversation"
    item = {
        "source": "whatsapp",
        "source_id": str(path),
        "kind": "whatsapp-conversation",
        "title": meta.get("summary") or first_line,
        "author": meta.get("peopleMentioned", [""])[0] if isinstance(meta.get("peopleMentioned"), list) and meta.get("peopleMentioned") else "",
        "recipients": "ExampleCo",
        "created_at": meta.get("taggedAt") or datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
        "raw_path": str(path),
        "checksum": checksum,
        "metadata_json": safe_json_dumps(meta),
    }
    return upsert_item(con, item, transcript, [])


def detect_source(path):
    p = Path(path)
    norm = p.as_posix().lower()
    if "data/sessions/raw/claude-code/" in norm:
        return "claude-code"
    if "data/sessions/raw/codex/" in norm:
        return "codex"
    if "data/sessions/raw/cursor/" in norm or "/.cursor/" in norm:
        return "cursor"
    if "data/sessions/raw/chatgpt/" in norm:
        return "chatgpt"
    if "data/sessions/raw/gemini/" in norm:
        return "gemini"
    if norm.endswith("message.json") and "data/gmail/raw/" in norm:
        return "gmail-json"
    if p.suffix.lower() == ".eml":
        return "eml"
    if "data/otter/raw/" in norm and p.suffix.lower() == ".json":
        return "otter"
    if "data/vapi/raw/" in norm and p.suffix.lower() == ".json":
        return "vapi"
    if "data/linkedin/raw/" in norm and p.suffix.lower() == ".json":
        return "linkedin"
    if "data/whatsapp/raw/" in norm and p.suffix.lower() == ".json":
        return "whatsapp-json"
    if "data/whatsapp/raw/" in norm and p.suffix.lower() == ".txt":
        return "whatsapp-conversation"
    if "data/sms/raw/" in norm or "data/imessage/raw/" in norm:
        return "sms"
    if p.name == "linkedin-dms.json":
        return "linkedin"
    if p.name == "transcript.txt" and "conversation" in norm:
        return "whatsapp-conversation"
    if p.suffix.lower() == ".jsonl":
        return "jsonl"
    if p.suffix.lower() in (".json", ".txt", ".md"):
        return "generic"
    return "file"


def index_file(con, path, source=None, copy_external=False, source_root=None):
    path = Path(path)
    if not path.exists() or not path.is_file():
        return 0
    kind = source or detect_source(path)
    original_path = path
    if copy_external and not str(path.resolve()).lower().startswith(str(REPO.resolve()).lower()):
        path = archive_external_file(path, kind)
    try:
        if kind == "migrated-file":
            return index_migrated_file(con, path, source_root or path.parent)
        if kind == "gmail-json":
            return 1 if index_gmail_message_json(con, path) else 0
        if kind == "eml":
            return 1 if index_eml(con, path) else 0
        if kind == "otter":
            return 1 if index_otter_json(con, path) else 0
        if kind == "vapi":
            return 1 if index_vapi_json(con, path) else 0
        if kind == "linkedin":
            result = index_linkedin_json(con, path)
            return result if isinstance(result, int) else 1
        if kind == "whatsapp-conversation":
            return 1 if index_whatsapp_conversation(con, path) else 0
        if kind == "whatsapp-json":
            return 1 if index_structured_raw_file(con, path, "whatsapp", "conversation-metadata") else 0
        if kind in ("sms", "imessage"):
            if path.suffix.lower() not in SMS_ARCHIVE_EXTENSIONS:
                return 0
            checksum = sha256_file(path)
            if item_checksum_matches(con, item_id_for("sms", str(path), str(path)), checksum):
                return 0
            text = extract_attachment_text(path)
            if not text and path.suffix.lower() in (".csv", ".txt", ".vcf", ".url"):
                text = path.read_text(encoding="utf-8", errors="replace")
            if not text:
                text = path.name
            item = {
                "source": "sms",
                "source_id": str(path),
                "kind": "sms-export-file",
                "title": f"SMS export: {path.name}",
                "author": "",
                "recipients": "ExampleCo",
                "created_at": datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
                "raw_path": str(path),
                "checksum": checksum,
                "metadata_json": safe_json_dumps({"original_path": str(original_path), "extension": path.suffix.lower()}),
            }
            upsert_item(con, item, text, [])
            return 1
        if kind in ("claude-code", "codex", "cursor", "chatgpt", "gemini", "session"):
            return 1 if index_session_file(con, path, kind, copy_external=False) else 0
        if kind in ("jsonl", "dispatch"):
            return index_jsonl(con, path, "dispatch")
        if kind == "generic":
            checksum = sha256_file(path)
            if item_checksum_matches(con, item_id_for("generic", str(path), str(path)), checksum):
                return 0
            text = path.read_text(encoding="utf-8", errors="replace")
            item = {
                "source": "generic",
                "source_id": str(path),
                "kind": "file",
                "title": path.name,
                "author": "",
                "recipients": "",
                "created_at": datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
                "raw_path": str(path),
                "checksum": checksum,
                "metadata_json": safe_json_dumps({"original_path": str(original_path)}),
            }
            upsert_item(con, item, text, [])
            return 1
    except Exception as e:
        print(f"[life-archive] index failed {path}: {e}", file=sys.stderr)
        return 0
    return 0


def iter_files(root, patterns=None):
    root = Path(root)
    if not root.exists():
        return
    patterns = patterns or ["*"]
    for pat in patterns:
        yield from root.rglob(pat)


def index_path(con, path, source=None, copy_external=False, limit=0, max_seconds=0):
    path = Path(path)
    count = 0
    started = time.monotonic()
    if path.is_file():
        count += index_file(con, path, source=source, copy_external=copy_external)
    elif path.is_dir():
        norm_root = path.as_posix().lower()
        if "data/gmail/raw" in norm_root:
            candidates = [f for f in path.rglob("message.json") if f.is_file()]
            candidates.sort(key=lambda f: f.stat().st_mtime, reverse=True)
        else:
            candidates = path.rglob("*")
        for f in candidates:
            if f.is_file():
                if f.name.endswith((".tmp", ".lock", ".text.txt")):
                    continue
                norm = f.as_posix().lower()
                if "data/gmail/raw/" in norm and "/attachments/" in norm:
                    continue
                if "data/gmail/raw/" in norm and f.name.lower() == "raw.eml" and f.with_name("message.json").exists():
                    continue
                count += index_file(con, f, source=source, copy_external=copy_external)
                if count % 100 == 0:
                    con.commit()
                if limit and count >= limit:
                    break
                if max_seconds and (time.monotonic() - started) >= max_seconds:
                    break
    con.commit()
    return count


def write_json_atomic(path, value):
    path = Path(path)
    ensure_parent(path)
    tmp = path.with_suffix(path.suffix + ".tmp")
    tmp.write_text(safe_json_dumps(value, indent=2) + "\n", encoding="utf-8")
    os.replace(tmp, path)


def migrated_rows(con, source_root):
    wanted_root = str(Path(source_root).resolve()).lower()
    rows = []
    for row in con.execute(
        "SELECT id, source_id, title, created_at, raw_path, checksum, metadata_json, indexed_at "
        "FROM items WHERE source = 'migrated-file' ORDER BY source_id"
    ):
        try:
            meta = json.loads(row["metadata_json"] or "{}")
        except Exception:
            meta = {}
        if str(meta.get("migration_root") or "").lower() != wanted_root:
            continue
        rows.append((row, meta))
    return rows


def write_migrated_manifest(con, source_root, manifest_path):
    manifest_path = Path(manifest_path)
    ensure_parent(manifest_path)
    tmp = manifest_path.with_suffix(manifest_path.suffix + ".tmp")
    with tmp.open("w", encoding="utf-8", newline="\n") as out:
        for row, meta in migrated_rows(con, source_root):
            record = {
                "item_id": row["id"],
                "absolute_path": row["raw_path"],
                "source_relative_path": meta.get("source_relative_path") or row["source_id"],
                "filename": meta.get("filename") or row["title"],
                "extension": meta.get("extension"),
                "mime_type": meta.get("mime_type"),
                "size": meta.get("size"),
                "sha256": row["checksum"],
                "created_at": row["created_at"],
                "filesystem_created": meta.get("filesystem_created"),
                "filesystem_modified": meta.get("filesystem_modified"),
                "embedded_created": meta.get("embedded_created"),
                "embedded_modified": meta.get("embedded_modified"),
                "filename_date": meta.get("filename_date"),
                "timestamp_provenance": meta.get("timestamp_provenance"),
                "extraction_status": meta.get("extraction_status"),
                "extracted_text_chars": meta.get("extracted_text_chars", 0),
                "sensitivity": meta.get("sensitivity"),
                "graphiti_eligible": bool(meta.get("graphiti_eligible")),
                "indexed_at": row["indexed_at"],
            }
            out.write(safe_json_dumps(record) + "\n")
    os.replace(tmp, manifest_path)


def migrated_coverage(con, source_root, scanned_files, scanned_bytes, errors, complete):
    rows = migrated_rows(con, source_root)
    status_counts = {}
    provenance_counts = {}
    indexed_bytes = 0
    graphiti_eligible = 0
    sensitive = 0
    extracted_items = 0
    extracted_chars = 0
    for _, meta in rows:
        status = meta.get("extraction_status") or "ExampleCo"
        status_counts[status] = status_counts.get(status, 0) + 1
        provenance = meta.get("timestamp_provenance") or "ExampleCo"
        provenance_counts[provenance] = provenance_counts.get(provenance, 0) + 1
        indexed_bytes += int(meta.get("size") or 0)
        graphiti_eligible += 1 if meta.get("graphiti_eligible") else 0
        sensitive += 1 if meta.get("sensitivity") else 0
        chars = int(meta.get("extracted_text_chars") or 0)
        if chars:
            extracted_items += 1
            extracted_chars += chars
    return {
        "schema_version": 1,
        "generated_at": now_iso(),
        "source_root": str(Path(source_root).resolve()),
        "complete": bool(complete),
        "scanned_files": scanned_files,
        "scanned_bytes": scanned_bytes,
        "indexed_items": len(rows),
        "indexed_bytes": indexed_bytes,
        "coverage_matches_scan": len(rows) == scanned_files and indexed_bytes == scanned_bytes,
        "sha256_coverage": len(rows),
        "extracted_text_items": extracted_items,
        "extracted_text_chars": extracted_chars,
        "metadata_only_items": len(rows) - extracted_items,
        "sensitive_metadata_only_items": sensitive,
        "graphiti_eligible_items": graphiti_eligible,
        "extraction_status_counts": status_counts,
        "timestamp_provenance_counts": provenance_counts,
        "ocr_available": bool(tesseract_path()),
        "errors": errors,
        "incremental_refresh": (
            "Rerun index-migrated with the same root, database, manifest, and receipt. "
            "Files with unchanged size and nanosecond mtime retain their prior SHA-256 and extraction; "
            "changed files are re-hashed and re-extracted."
        ),
        "timestamp_limitation": (
            "The AWS transfer did not preserve original filesystem mtimes. Embedded metadata and "
            "filename dates are recorded when available; otherwise the filesystem time reflects migration."
        ),
    }


def index_migrated_path(con, source_root, manifest_path, receipt_path, limit=0, max_seconds=0, workers=1):
    source_root = Path(source_root).resolve()
    if not source_root.exists() or not source_root.is_dir():
        raise ValueError(f"Migration root is missing or is not a directory: {source_root}")
    started = time.monotonic()
    scanned_files = 0
    scanned_bytes = 0
    changed = 0
    unchanged = 0
    pruned_stale = 0
    errors = []
    complete = True
    seen_relative_paths = set()
    existing = {
        (meta.get("source_relative_path") or row["source_id"]): meta
        for row, meta in migrated_rows(con, source_root)
    }
    candidates = sorted((p for p in source_root.rglob("*") if p.is_file()), key=lambda p: p.as_posix().lower())
    jobs = []
    for path in candidates:
        try:
            stat = path.stat()
            relative_path = path.relative_to(source_root).as_posix()
            seen_relative_paths.add(relative_path)
            scanned_files += 1
            scanned_bytes += stat.st_size
            prior = existing.get(relative_path) or {}
            if (
                int(prior.get("size") or -1) == stat.st_size
                and int(prior.get("filesystem_mtime_ns") or -1) == stat.st_mtime_ns
                and prior.get("sha256")
            ):
                unchanged += 1
            else:
                jobs.append((str(path), str(source_root)))
        except Exception as exc:
            errors.append({"path": str(path), "error": str(exc)[:500]})
        if limit and scanned_files >= limit:
            complete = False
            break
        if max_seconds and (time.monotonic() - started) >= max_seconds:
            complete = False
            break
    worker_count = max(1, int(workers or 1))
    if max_seconds:
        # A bounded partial run must be able to stop promptly; queued process-pool
        # work would violate that contract, so keep it serial.
        worker_count = 1
    if worker_count > 1 and jobs:
        with ProcessPoolExecutor(max_workers=worker_count) as pool:
            prepared_results = pool.map(prepare_migrated_file_safe, jobs, chunksize=1)
            for completed, result in enumerate(prepared_results, start=1):
                if result.get("ok"):
                    changed += store_prepared_migrated_file(con, result["prepared"])
                else:
                    errors.append({"path": result.get("path"), "error": result.get("error")})
                if completed % 100 == 0:
                    con.commit()
    else:
        for completed, job in enumerate(jobs, start=1):
            result = prepare_migrated_file_safe(job)
            if result.get("ok"):
                changed += store_prepared_migrated_file(con, result["prepared"])
            else:
                errors.append({"path": result.get("path"), "error": result.get("error")})
            if completed % 100 == 0:
                con.commit()
            if max_seconds and (time.monotonic() - started) >= max_seconds:
                complete = False
                break
    con.commit()
    if complete:
        for row, meta in migrated_rows(con, source_root):
            relative_path = meta.get("source_relative_path") or row["source_id"]
            if relative_path in seen_relative_paths:
                continue
            con.execute("DELETE FROM attachments WHERE item_id = ?", (row["id"],))
            con.execute("DELETE FROM item_fts WHERE item_id = ?", (row["id"],))
            con.execute("DELETE FROM items WHERE id = ?", (row["id"],))
            pruned_stale += 1
        con.commit()
        # Re-scan totals after indexing so the receipt proves the exact final source view.
        final_files = [p for p in source_root.rglob("*") if p.is_file()]
        scanned_files = len(final_files)
        scanned_bytes = sum(p.stat().st_size for p in final_files)
    write_migrated_manifest(con, source_root, manifest_path)
    receipt = migrated_coverage(con, source_root, scanned_files, scanned_bytes, errors, complete)
    receipt["changed_or_new_items"] = changed
    receipt["unchanged_items_fast_path"] = unchanged
    receipt["pruned_stale_index_items"] = pruned_stale
    receipt["extraction_workers"] = worker_count
    receipt["database"] = str(Path(con.execute("PRAGMA database_list").fetchone()[2]).resolve())
    receipt["manifest"] = str(Path(manifest_path).resolve())
    write_json_atomic(receipt_path, receipt)
    return receipt


def default_existing_paths():
    paths = [
        (REPO / "data" / "gmail" / "raw", None, False),
        (REPO / "data" / "otter" / "raw", None, False),
        (REPO / "data" / "vapi" / "raw", None, False),
        (REPO / "data" / "linkedin" / "raw", None, False),
        (REPO / "data" / "sms" / "raw", "sms", False),
        (REPO / "data" / "imessage" / "raw", "imessage", False),
        (REPO / "data" / "whatsapp" / "raw", None, False),
        (REPO / "data" / "sessions" / "raw" / "claude-code", "claude-code", False),
        (REPO / "data" / "sessions" / "raw" / "claude-code-mirror", "claude-code", False),
        (REPO / "data" / "sessions" / "raw" / "codex", "codex", False),
        (REPO / "data" / "sessions" / "raw" / "cursor", "cursor", False),
        (REPO / "data" / "sessions" / "raw" / "chatgpt", "chatgpt", False),
        (REPO / "data" / "sessions" / "raw" / "gemini", "gemini", False),
        (REPO / "data" / "agent" / "dispatch-queue.jsonl", "dispatch", False),
        (REPO / "data" / "agent" / "amy-dispatch-log.jsonl", "dispatch", False),
        (REPO / "data" / "agent" / "dispatch-resolutions.jsonl", "dispatch", False),
        (APPDATA / "secondbrain" / "data" / "agent" / "linkedin-dms.json", None, True),
        (APPDATA / "secondbrain" / "data" / "conversations", None, True),
        (APPDATA / "secondbrain" / "data" / "calls", "vapi", True),
        (Path.home() / ".claude" / "projects", "claude-code", True),
        (Path.home() / ".codex" / "sessions", "codex", True),
        (Path.home() / ".cursor" / "projects", "cursor", True),
    ]
    return paths


def index_existing(con, include_external=True, limit=0, max_seconds=0):
    total = 0
    started = time.monotonic()
    for path, source, external in default_existing_paths():
        if external and not include_external:
            continue
        remaining_limit = max(0, limit - total) if limit else 0
        remaining_seconds = max(1, max_seconds - (time.monotonic() - started)) if max_seconds else 0
        total += index_path(con, path, source=source, copy_external=external, limit=remaining_limit, max_seconds=remaining_seconds)
        if limit and total >= limit:
            break
        if max_seconds and (time.monotonic() - started) >= max_seconds:
            break
    con.commit()
    return total


def normalize_fts_query(query):
    query = str(query or "").strip()
    if not query:
        return '""'
    if re.search(r'["*:(){}+\-^]', query):
        terms = re.findall(r"[A-Za-z0-9_@.]+", query)
        return " ".join(terms) if terms else '"' + query.replace('"', '""') + '"'
    return query


def search(con, query, limit=20, source=None):
    fts_query = normalize_fts_query(query)
    source_filter = source or None
    rows = con.execute(
        """
        SELECT i.*, bm25(item_fts) AS rank,
               snippet(item_fts, 3, '[', ']', ' ... ', 24) AS body_snippet,
               snippet(item_fts, 4, '[', ']', ' ... ', 24) AS attachment_snippet
        FROM item_fts
        JOIN items i ON i.id = item_fts.item_id
        WHERE item_fts MATCH ?
          AND (? IS NULL OR i.source = ?)
        ORDER BY rank
        LIMIT ?
        """,
        (fts_query, source_filter, source_filter, limit),
    ).fetchall()
    return [dict(r) for r in rows]


def gmail_hex_ids(meta_bytes):
    msgid_match = re.search(rb"X-GM-MSGID (\d+)", meta_bytes or b"")
    thrid_match = re.search(rb"X-GM-THRID (\d+)", meta_bytes or b"")
    gm_msg_hex = format(int(msgid_match.group(1)), "x") if msgid_match else None
    gm_thr_hex = format(int(thrid_match.group(1)), "x") if thrid_match else None
    return gm_msg_hex, gm_thr_hex


def gmail_fetch_meta_bytes(fetched):
    parts = []
    for item in fetched or []:
        if isinstance(item, tuple):
            if isinstance(item[0], bytes):
                parts.append(item[0])
            else:
                parts.append(str(item[0]).encode())
        elif isinstance(item, bytes):
            parts.append(item)
    return b" ".join(parts)


def gmail_fetch_raw_batch(client, eids):
    if not eids:
        return []
    request_ids = b",".join(eids).decode("ascii", errors="ignore")
    typ, fetched = client.fetch(request_ids, "(X-GM-THRID X-GM-MSGID RFC822)")
    if typ != "OK":
        return []
    rows = []
    fallback = iter(eids)
    for item in fetched or []:
        if not isinstance(item, tuple):
            continue
        meta = item[0] if isinstance(item[0], bytes) else str(item[0]).encode()
        raw = item[1]
        m = re.match(rb"\s*(\d+)\s+", meta)
        eid = m.group(1) if m else next(fallback, b"")
        if raw:
            rows.append((eid, meta, raw))
    return rows


def load_existing_gmail_seen(raw_dir):
    seen = set()
    root = Path(raw_dir)
    if not root.exists():
        return seen
    for msg_json in root.rglob("message.json"):
        try:
            item = json.loads(msg_json.read_text(encoding="utf-8", errors="replace"))
        except Exception:
            continue
        gm_msg_hex = item.get("gmail_message_hex")
        if gm_msg_hex:
            seen.add(str(gm_msg_hex))
    return seen


def save_gmail_message(root, eid, raw_bytes, meta_bytes, extract_attachments=True, include_body=True):
    msg = email.message_from_bytes(raw_bytes)
    message_id = (msg.get("Message-ID") or "").strip("<>") or eid.decode()
    date_hdr = decode_hdr(msg.get("Date"))
    if parse_email_date(date_hdr) is None:
        date_hdr = received_date_header(msg)
    dt = parse_email_date(date_hdr)
    y, m, d = date_parts(dt)
    gm_msg_hex, gm_thr_hex = gmail_hex_ids(meta_bytes or b"")
    stable = gm_msg_hex or message_id
    out_dir = Path(root) / y / m / d / safe_segment(stable)
    out_dir.mkdir(parents=True, exist_ok=True)
    raw_path = out_dir / "raw.eml"
    raw_path.write_bytes(raw_bytes)
    attachments = []
    if extract_attachments:
        for idx, part in enumerate(msg.walk()):
            filename = part.get_filename()
            disp = str(part.get("Content-Disposition") or "")
            cid = (part.get("Content-ID") or "").strip("<>")
            if not (filename or cid or "attachment" in disp.lower() or "inline" in disp.lower()):
                continue
            payload = part.get_payload(decode=True) or b""
            if not payload:
                continue
            att_dir = out_dir / "attachments"
            att_dir.mkdir(exist_ok=True)
            name = safe_segment(decode_hdr(filename or cid or f"part-{idx}"))
            if "." not in name and part.get_content_subtype():
                name += "." + part.get_content_subtype()
            att_path = att_dir / name
            att_path.write_bytes(payload)
            text = extract_attachment_text(att_path, part.get_content_type())
            attachments.append({
                "filename": decode_hdr(filename or cid or name),
                "content_type": part.get_content_type(),
                "size": len(payload),
                "content_id": cid or None,
                "disposition": disp or None,
                "archive_saved_to": str(att_path),
                "text_saved_to": str(att_path.with_suffix(att_path.suffix + ".text.txt")) if text else None,
            })
    item = {
        "id": eid.decode() if isinstance(eid, bytes) else str(eid),
        "message_id": message_id,
        "gmail_message_hex": gm_msg_hex,
        "thread_hex": gm_thr_hex,
        "gmail_url": f"https://mail.google.com/mail/u/0/#all/{gm_thr_hex}" if gm_thr_hex else None,
        "from": decode_hdr(msg.get("From")),
        "to": decode_hdr(msg.get("To")),
        "cc": decode_hdr(msg.get("Cc")),
        "subject": decode_hdr(msg.get("Subject")),
        "date": date_hdr,
        "body": body_from_email_message(msg) if include_body else "",
        "attachments": attachments,
        "attachments_embedded_in_raw_eml": not extract_attachments,
        "archive_dir": str(out_dir),
        "raw_saved_to": str(raw_path),
        "archived_at": now_iso(),
    }
    (out_dir / "message.json").write_text(safe_json_dumps(item, indent=2), encoding="utf-8")
    return out_dir / "message.json"


def gmail_backfill(args):
    user = (SECRETS / "gmail_sender.txt").read_text().strip()
    password = (SECRETS / "gmail_app_password.txt").read_text().strip()
    con = None if args.defer_index else connect_db(args.db)
    state_path = Path(args.state_file) if args.state_file and not args.no_state else None
    state = {"seen": [], "updatedAt": None}
    seen = set()
    if state_path and state_path.exists():
        try:
            state = json.loads(state_path.read_text(encoding="utf-8", errors="replace"))
            seen = set(state.get("seen") or [])
        except Exception:
            state = {"seen": [], "updatedAt": None}
            seen = set()
    local_seen = set() if args.skip_local_seen_scan else load_existing_gmail_seen(args.raw_dir)
    seen.update(local_seen)
    started = time.monotonic()
    client = imaplib.IMAP4_SSL("imap.gmail.com", timeout=45)
    client.login(user, password)
    try:
        client.select(f'"{args.mailbox}"', readonly=True)
        if args.query:
            q = '"' + args.query.replace('"', '\\"') + '"'
            typ, data = client.search(None, "X-GM-RAW", q)
        else:
            typ, data = client.search(None, "ALL")
        ids = data[0].split() if data and data[0] else []
        if args.newest_first:
            ids = list(reversed(ids))
        start_index = 0
        if not args.restart_scan:
            try:
                if state.get("mailbox") == args.mailbox and state.get("query") == args.query:
                    start_index = max(0, min(int(state.get("lastScanIndex") or 0), len(ids)))
            except Exception:
                start_index = 0
        processed = 0
        skipped = 0
        scanned = start_index
        stopped_reason = None

        def write_state():
            if not state_path:
                return
            state_path.parent.mkdir(parents=True, exist_ok=True)
            state_path.write_text(safe_json_dumps({
                "seen": sorted(seen),
                "localSeeded": len(local_seen),
                "lastMatched": len(ids),
                "lastScanIndex": scanned,
                "lastProcessed": processed,
                "lastSkipped": skipped,
                "mailbox": args.mailbox,
                "query": args.query,
                "newestFirst": bool(args.newest_first),
                "complete": scanned >= len(ids),
                "stoppedReason": stopped_reason,
                "updatedAt": now_iso(),
            }, indent=2), encoding="utf-8")

        batch_size = max(1, int(args.batch_size or 1))
        state_flush_interval = max(1, int(getattr(args, "state_flush_interval", 25) or 25))
        remaining_ids = ids[start_index:]
        for batch_start in range(0, len(remaining_ids), batch_size):
            if args.max_seconds and time.monotonic() - started >= args.max_seconds:
                stopped_reason = "max_seconds"
                break
            batch = remaining_ids[batch_start:batch_start + batch_size]
            fetched_rows = gmail_fetch_raw_batch(client, batch)
            fetched_by_eid = {eid: (meta, raw) for eid, meta, raw in fetched_rows}
            for eid in batch:
                if args.max_seconds and time.monotonic() - started >= args.max_seconds:
                    stopped_reason = "max_seconds"
                    break
                scanned += 1
                meta, raw = fetched_by_eid.get(eid, (b"", None))
                if not raw:
                    typ, fetched = client.fetch(eid, "(X-GM-THRID X-GM-MSGID RFC822)")
                    raw = None
                    meta = b""
                    for item in fetched:
                        if isinstance(item, tuple):
                            meta += item[0] if isinstance(item[0], bytes) else str(item[0]).encode()
                            raw = item[1]
                gm_msg_hex, _ = gmail_hex_ids(meta)
                if gm_msg_hex and gm_msg_hex in seen:
                    skipped += 1
                    continue
                if not raw:
                    skipped += 1
                    continue
                msg_json = save_gmail_message(
                    args.raw_dir,
                    eid,
                    raw,
                    meta,
                    extract_attachments=not args.raw_only,
                    include_body=not args.defer_body,
                )
                if not args.defer_index:
                    index_gmail_message_json(con, msg_json)
                if gm_msg_hex:
                    seen.add(gm_msg_hex)
                processed += 1
                if (processed + skipped) % state_flush_interval == 0:
                    if con is not None:
                        con.commit()
                    write_state()
                    print(f"[life-archive] gmail backfill scanned={scanned}/{len(ids)} processed={processed} skipped={skipped}", file=sys.stderr)
                if args.limit and processed >= args.limit:
                    stopped_reason = "limit"
                    break
            if stopped_reason:
                break
            if len(fetched_rows) < len(batch):
                skipped += 1
        if con is not None:
            con.commit()
        write_state()
        print(json.dumps({"processed": processed, "skipped": skipped, "scanned": scanned, "matched": len(ids), "complete": scanned >= len(ids), "stateFile": str(state_path) if state_path else None}, indent=2))
    finally:
        client.logout()


def print_stats(con):
    by_source = con.execute("SELECT source, COUNT(*) c FROM items GROUP BY source ORDER BY c DESC").fetchall()
    attachments = con.execute("SELECT COUNT(*) c, SUM(CASE WHEN extracted_text IS NOT NULL AND length(extracted_text) > 0 THEN 1 ELSE 0 END) t FROM attachments").fetchone()
    print(json.dumps({
        "items_by_source": [dict(r) for r in by_source],
        "attachments": dict(attachments),
    }, indent=2, ensure_ascii=False))


def print_item(con, item_id):
    row = con.execute("SELECT * FROM items WHERE id = ? LIMIT 1", (item_id,)).fetchone()
    print(safe_json_dumps(dict(row) if row else None, indent=2))


def compact_stale(con):
    rows = con.execute("SELECT id, source, raw_path FROM items").fetchall()
    stale = []
    seen_sms_paths = set()
    for row in rows:
        raw_path = row["raw_path"]
        if not raw_path:
            continue
        p = Path(raw_path)
        remove = False
        if row["source"] in {"sms", "imessage"}:
            remove = (not p.exists()) or p.suffix.lower() not in SMS_ARCHIVE_EXTENSIONS
            try:
                canonical = str(p.resolve()).lower()
            except Exception:
                canonical = str(p).lower()
            if not remove and canonical in seen_sms_paths:
                remove = True
            seen_sms_paths.add(canonical)
        elif str(p).lower().startswith(str(REPO).lower()) and not p.exists():
            remove = True
        if remove:
            stale.append(row["id"])
    if not stale:
        print(json.dumps({"removed": 0}, indent=2))
        return
    fts_rowids = {
        row["item_id"]: row["rowid"]
        for row in con.execute("SELECT rowid, item_id FROM item_fts").fetchall()
    }
    for item_id in stale:
        con.execute("DELETE FROM attachments WHERE item_id = ?", (item_id,))
        rowid = fts_rowids.get(item_id)
        if rowid is not None:
            con.execute("DELETE FROM item_fts WHERE rowid = ?", (rowid,))
        con.execute("DELETE FROM items WHERE id = ?", (item_id,))
    con.commit()
    print(json.dumps({"removed": len(stale)}, indent=2))


def main():
    parser = argparse.ArgumentParser(description="Build and search the SecondBrain life archive.")
    parser.add_argument("--db", default=str(DEFAULT_DB))
    sub = parser.add_subparsers(dest="cmd", required=True)

    p = sub.add_parser("index-path")
    p.add_argument("--path", required=True)
    p.add_argument("--source")
    p.add_argument("--copy-external", action="store_true")
    p.add_argument("--limit", type=int, default=0, help="Maximum new or changed items to index.")
    p.add_argument("--max-seconds", type=int, default=0, help="Stop after this many seconds.")
    p.add_argument("--wait-lock-seconds", type=int, default=0, help="Wait up to this long for another archive writer to finish.")

    p = sub.add_parser("index-existing")
    p.add_argument("--no-external", action="store_true")
    p.add_argument("--limit", type=int, default=0, help="Maximum new or changed items to index.")
    p.add_argument("--max-seconds", type=int, default=0, help="Stop after this many seconds.")
    p.add_argument("--wait-lock-seconds", type=int, default=0, help="Wait up to this long for another archive writer to finish.")

    p = sub.add_parser("index-migrated")
    p.add_argument("--path", required=True, help="Root of the landed migration tree.")
    p.add_argument("--manifest", required=True, help="Rebuildable JSONL file manifest output.")
    p.add_argument("--receipt", required=True, help="Coverage receipt JSON output.")
    p.add_argument("--limit", type=int, default=0, help="Maximum files to scan in this run.")
    p.add_argument("--max-seconds", type=int, default=0, help="Stop after this many seconds and write a partial receipt.")
    p.add_argument("--workers", type=int, default=min(6, os.cpu_count() or 1), help="Parallel extraction worker processes.")
    p.add_argument("--wait-lock-seconds", type=int, default=0, help="Wait up to this long for another archive writer to finish.")

    p = sub.add_parser("search")
    p.add_argument("query")
    p.add_argument("--limit", type=int, default=20)
    p.add_argument("--source", help="Filter to one source, e.g. gmail, otter, linkedin, whatsapp, claude-code, codex.")

    p = sub.add_parser("stats")

    p = sub.add_parser("item")
    p.add_argument("item_id")

    p = sub.add_parser("compact-stale")
    p.add_argument("--wait-lock-seconds", type=int, default=0, help="Wait up to this long for another archive writer to finish.")

    p = sub.add_parser("backfill-gmail")
    p.add_argument("--mailbox", default="[Gmail]/All Mail")
    p.add_argument("--raw-dir", default=str(DEFAULT_GMAIL_RAW))
    p.add_argument("--query", help='Gmail X-GM-RAW query, e.g. newer_than:30d')
    p.add_argument("--limit", type=int, default=0)
    p.add_argument("--max-seconds", type=int, default=0, help="Stop after this many seconds.")
    p.add_argument("--batch-size", type=int, default=25, help="Gmail messages to fetch per IMAP request.")
    p.add_argument("--raw-only", action="store_true", help="Save raw .eml and index body first; defer attachment extraction/OCR.")
    p.add_argument("--defer-index", action="store_true", help="Save raw Gmail files and state first; leave SQLite indexing for index-existing.")
    p.add_argument("--defer-body", action="store_true", help="Save raw Gmail bytes and metadata first; parse body later during indexing.")
    p.add_argument("--newest-first", action="store_true")
    p.add_argument("--restart-scan", action="store_true", help="Ignore the saved Gmail scan cursor and rescan from the beginning.")
    p.add_argument("--skip-local-seen-scan", action="store_true", help="Trust the saved Gmail state seen set instead of rescanning local raw files before backfill.")
    p.add_argument("--state-flush-interval", type=int, default=25, help="Write Gmail backfill state every N scanned messages, plus once at the end.")
    p.add_argument("--wait-lock-seconds", type=int, default=0, help="Wait up to this long for another archive writer to finish.")
    p.add_argument("--state-file", default=str(DEFAULT_GMAIL_BACKFILL_STATE))
    p.add_argument("--no-state", action="store_true")

    args = parser.parse_args()
    lock_fd = None
    lock_path = lock_path_for_db(args.db)
    mutating = args.cmd in ("backfill-gmail", "index-path", "index-existing", "index-migrated", "compact-stale")
    if mutating:
        deadline = time.monotonic() + max(0, int(getattr(args, "wait_lock_seconds", 0) or 0))
        while True:
            lock_fd = acquire_lock(lock_path)
            if lock_fd is not None:
                break
            if time.monotonic() >= deadline:
                print(json.dumps({"skipped": True, "reason": "life archive index/backfill already running", "lock": str(lock_path)}, indent=2))
                return
            time.sleep(5)
    try:
        if args.cmd == "backfill-gmail":
            gmail_backfill(args)
            return
        con = connect_db(args.db)
        if args.cmd == "index-path":
            count = index_path(con, args.path, source=args.source, copy_external=args.copy_external, limit=args.limit, max_seconds=args.max_seconds)
            print(json.dumps({"indexed": count, "db": args.db}, indent=2))
        elif args.cmd == "index-migrated":
            receipt = index_migrated_path(
                con, args.path, args.manifest, args.receipt,
                limit=args.limit, max_seconds=args.max_seconds, workers=args.workers,
            )
            print(safe_json_dumps(receipt, indent=2))
        elif args.cmd == "index-existing":
            count = index_existing(con, include_external=not args.no_external, limit=args.limit, max_seconds=args.max_seconds)
            print(json.dumps({"indexed": count, "db": args.db}, indent=2))
        elif args.cmd == "search":
            print(safe_json_dumps(search(con, args.query, args.limit, source=args.source), indent=2))
        elif args.cmd == "stats":
            print_stats(con)
        elif args.cmd == "item":
            print_item(con, args.item_id)
        elif args.cmd == "compact-stale":
            compact_stale(con)
    finally:
        release_lock(lock_fd, lock_path)


if __name__ == "__main__":
    main()
